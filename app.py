import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from PyPDF2 import PdfReader
from pptx import Presentation

st.set_page_config(page_title="CredIntel AI", layout="wide")

st.title("CredIntel AI - Credit Intelligence System")

st.markdown("Flexible data extraction system for corporate credit analysis")

# --------------------------------------------------
# USER CHOOSES INPUT TYPE
# --------------------------------------------------

input_type = st.radio(
"Select Data Input Method",
["Paste Text Data","Upload File"]
)

# --------------------------------------------------
# TEXT INPUT MODE
# --------------------------------------------------

text_data = ""

if input_type == "Paste Text Data":

    text_data = st.text_area(
    "Paste all company data here (financial statements, annual report insights, legal issues, sector news etc.)",
    height=300
    )

# --------------------------------------------------
# FILE INPUT MODE
# --------------------------------------------------

if input_type == "Upload File":

    uploaded_file = st.file_uploader(
    "Upload company document (PDF, PPTX, Excel, TXT, CSV)",
    type=["pdf","pptx","xlsx","xls","txt","csv"]
    )

    def extract_text_from_file(file):

        if file is None:
            return ""

        file_type = file.name.split(".")[-1]

        extracted = ""

        # PDF
        if file_type == "pdf":

            reader = PdfReader(file)

            for page in reader.pages:
                extracted += page.extract_text()

        # PowerPoint
        elif file_type == "pptx":

            prs = Presentation(file)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape,"text"):
                        extracted += shape.text + " "

        # Excel
        elif file_type in ["xlsx","xls"]:

            df = pd.read_excel(file)

            extracted = df.to_string()

        # CSV
        elif file_type == "csv":

            df = pd.read_csv(file)

            extracted = df.to_string()

        # TXT
        elif file_type == "txt":

            extracted = str(file.read(),"utf-8")

        return extracted.lower()

    text_data = extract_text_from_file(uploaded_file)

# --------------------------------------------------
# EXTRA USER DETAILS
# --------------------------------------------------

st.header("Additional Company Details")

company_name = st.text_input("Company Name")

sector = st.text_input("Sector")

loan_amount = st.number_input("Requested Loan Amount")

revenue = st.number_input("Annual Revenue", value=0)

profit = st.number_input("Net Profit / Loss", value=0)

debt = st.number_input("Total Debt Outstanding", value=0)

gst_turnover = st.number_input("GST Turnover", value=0)

# --------------------------------------------------
# CREDIT SCORING ENGINE
# --------------------------------------------------

def calculate_credit_score(data):

    score = 80

    reasons = []

    # Financial checks

    if profit < 0:
        score -= 20
        reasons.append("Company reporting losses")

    if debt > revenue and revenue != 0:
        score -= 15
        reasons.append("Debt higher than revenue")

    if gst_turnover == 0:
        score -= 10
        reasons.append("GST turnover unavailable")

    # Keyword analysis

    risk_keywords = {
        "litigation":20,
        "fraud":25,
        "default":25,
        "loss":15,
        "debt":10,
        "slowdown":10,
        "shutdown":20,
        "bankruptcy":30,
        "legal notice":20
    }

    for word, penalty in risk_keywords.items():

        if word in data:
            score -= penalty
            reasons.append(f"Risk keyword detected: {word}")

    return score, reasons

# --------------------------------------------------
# GENERATE DECISION
# --------------------------------------------------

if st.button("Run Credit Analysis"):

    score, reasons = calculate_credit_score(text_data)

    if score >= 75:
        decision = "APPROVE"
        interest = "Base + 1%"
    elif score >= 55:
        decision = "REVIEW"
        interest = "Base + 3%"
    else:
        decision = "REJECT"
        interest = "High Risk Premium"

    st.header("Credit Analysis Result")

    col1,col2,col3 = st.columns(3)

    col1.metric("Credit Score",score)
    col2.metric("Loan Decision",decision)
    col3.metric("Interest Rate Suggestion",interest)

    st.subheader("Detailed Reasoning")

    if reasons:
        for r in reasons:
            st.write("-",r)
    else:
        st.write("No major risk indicators detected")

# --------------------------------------------------
# RISK RADAR
# --------------------------------------------------

    financial_risk = min(100,debt/(revenue+1)*50)
    legal_risk = 80 if "litigation" in text_data else 20
    sector_risk = 60 if "slowdown" in text_data else 30
    operational_risk = 60 if "capacity" in text_data else 30

    labels = ["Financial","Legal","Sector","Operational"]

    values = [financial_risk,legal_risk,sector_risk,operational_risk]

    fig = plt.figure(figsize=(5,5))
    ax = fig.add_subplot(111,polar=True)

    angles = [0,1.57,3.14,4.71]
    values_cycle = values + [values[0]]
    angles_cycle = angles + [angles[0]]

    ax.plot(angles_cycle,values_cycle)
    ax.fill(angles_cycle,values_cycle,alpha=0.2)

    ax.set_xticks(angles)
    ax.set_xticklabels(labels)

    st.pyplot(fig)

# --------------------------------------------------
# CREDIT REPORT
# --------------------------------------------------

    st.subheader("Credit Appraisal Memo")

    st.write(f"""
Company: {company_name}

Sector: {sector}

Requested Loan: {loan_amount}

Credit Score: {score}

Decision: {decision}

Suggested Interest Rate: {interest}

Key Risk Signals:

{reasons}

Conclusion:
Based on the financial indicators and extracted document insights, the system generated the above credit recommendation.
""")
